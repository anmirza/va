#!/usr/bin/env python
"""Extract all text from PowerPoint presentation - Complete Version"""

import sys
import subprocess

# Try to import python-pptx, install if needed
try:
    from pptx import Presentation
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

# Path to the PowerPoint file
pptx_file = r'c:\Users\zainy\OneDrive\Desktop\FREELANCE\requsitiiiii-latest\requisti-version\Feedback Round 2.pptx'

try:
    print(f"Opening: {pptx_file}\n")
    prs = Presentation(pptx_file)
    
    print(f"Total slides: {len(prs.slides)}\n")
    print("=" * 80)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n=== SLIDE {i+1} ===\n")
        text_found = False
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                print(shape.text)
                text_found = True
        
        if not text_found:
            print("[No text content found on this slide]")
        
        print("\n" + "=" * 80)
    
    print("\n✓ Extraction complete!")
    
except Exception as e:
    print(f"Error: {e}")
    sys.exit(1)
