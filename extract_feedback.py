import sys
import os
import subprocess

def install_pkg(pkg):
    subprocess.check_call([sys.executable, "-m", "pip", "install", pkg])

try:
    from pptx import Presentation
except ImportError:
    install_pkg('python-pptx')
    from pptx import Presentation

try:
    import pandas as pd
except ImportError:
    install_pkg('pandas')
    install_pkg('openpyxl')
    import pandas as pd

pptx_path = r'c:\Users\zainy\OneDrive\Desktop\FREELANCE\requisti-v0\Creative Feedback_v2.pptx'
agency_xlsx = r'c:\Users\zainy\OneDrive\Desktop\FREELANCE\requisti-v0\RFI Agency-updated.xlsx'
prod_xlsx = r'c:\Users\zainy\OneDrive\Desktop\FREELANCE\requisti-v0\RFI Production House-updated.xlsx'

with open('feedback_content.txt', 'w', encoding='utf-8') as f:
    f.write('=== PPTX CONTENT ===\n')
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            f.write(f'--- Slide {i+1} ---\n')
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + '\n')
    except Exception as e:
        f.write(f'Error reading PPTX: {e}\n')

    f.write('\n\n=== RFI Agency XLSX ===\n')
    try:
        df_agency = pd.read_excel(agency_xlsx, sheet_name=None)
        for sheet, df in df_agency.items():
            f.write(f'Sheet: {sheet}\n')
            f.write(df.head(50).to_string() + '\n\n')
    except Exception as e:
        f.write(f'Error reading Agency XLSX: {e}\n')

    f.write('\n\n=== RFI Prod House XLSX ===\n')
    try:
        df_prod = pd.read_excel(prod_xlsx, sheet_name=None)
        for sheet, df in df_prod.items():
            f.write(f'Sheet: {sheet}\n')
            f.write(df.head(50).to_string() + '\n\n')
    except Exception as e:
        f.write(f'Error reading Prod XLSX: {e}\n')

print("Extraction complete.")
