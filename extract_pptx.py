#!/usr/bin/env python
"""Extract all text from PowerPoint presentation"""

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation

prs = Presentation(r'c:\Users\zainy\OneDrive\Desktop\FREELANCE\requsitiiiii-latest\requisti-version\Feedback Round 2.pptx')

for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text)
    print()
